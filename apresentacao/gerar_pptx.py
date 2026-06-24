# -*- coding: utf-8 -*-
"""Gera a apresentação ModaCRM AI (.pptx) seguindo o roteiro de demo."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paleta da marca (mesma do app) ──
INK        = RGBColor(0x2B, 0x23, 0x20)
INK_SOFT   = RGBColor(0x7D, 0x71, 0x6B)
CREAM      = RGBColor(0xFA, 0xF7, 0xF2)
ACCENT     = RGBColor(0xC2, 0x55, 0x2B)
ACCENT_SFT = RGBColor(0xF5, 0xE3, 0xDA)
OK         = RGBColor(0x2E, 0x7D, 0x32)
WARN       = RGBColor(0xB2, 0x6A, 0x00)
DANGER     = RGBColor(0xC6, 0x28, 0x28)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BD    = RGBColor(0xE8, 0xE0, 0xD8)

SERIF = "Georgia"
SANS  = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


SLIDES = []  # (slide, is_dark) — usado para rodapé/numeração


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    SLIDES.append((s, bg == INK))
    return s


def box(s, x, y, w, h):
    return s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame


def text(frame, runs, size=18, color=INK, font=SANS, bold=False, align=PP_ALIGN.LEFT,
         space_after=6, line=None, first=False):
    """runs: str ou lista de (txt, dict-override)."""
    p = frame.paragraphs[0] if first else frame.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line: p.line_spacing = line
    if isinstance(runs, str):
        runs = [(runs, {})]
    for txt, ov in runs:
        run = p.add_run(); run.text = txt
        f = run.font
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.name = ov.get("font", font)
        f.color.rgb = ov.get("color", color)
    return p


def rect(s, x, y, w, h, fill=WHITE, line=CARD_BD, line_w=1.0, round=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def bar(s, y=1.18, x=0.9, w=2.2):
    rect(s, x, y, w, 0.07, fill=ACCENT, line=None)


def kicker(s, txt, x=0.9, y=0.6):
    text(box(s, x, y, 8, 0.5), txt.upper(), size=13, color=ACCENT, bold=True, first=True)


def title(s, txt, x=0.9, y=0.78, size=34, w=11.5, color=INK):
    text(box(s, x, y, w, 1.1), txt, size=size, color=color, font=SERIF, bold=True, first=True)
    bar(s, y=y + 0.62, x=x)


# ───────────────────────── 1. Capa ─────────────────────────
s = slide(INK)
rect(s, 0, 0, 0.35, SH.inches, fill=ACCENT, line=None, round=False)
tf = box(s, 1.1, 2.35, 11, 2.6)
text(tf, [("Moda", {}), ("CRM", {"color": ACCENT}), (" AI", {})],
     size=60, color=WHITE, font=SERIF, bold=True, first=True)
text(tf, "A plataforma de vendas, relacionamento e estoque para lojas de moda",
     size=22, color=RGBColor(0xD8,0xCF,0xC8), space_after=4)
text(tf, [("O “Salesforce da Moda” — ", {"color": RGBColor(0xA8,0x9D,0x96)}),
          ("CRM + WhatsApp + IA por vendedora", {"color": RGBColor(0xE8,0xA8,0x7C), "bold": True})],
     size=18)
text(box(s, 1.1, 6.4, 11, 0.5), "Demonstração do produto  ·  Fases 1 a 3 concluídas",
     size=15, color=RGBColor(0xA8,0x9D,0x96), first=True)

# ───────────────────────── 2. O problema ─────────────────────────
s = slide()
kicker(s, "Por que existimos")
title(s, "A loja de moda perde venda no relacionamento")
dores = [
    ("🧩", "Carteira na cabeça da vendedora", "Quando ela sai, os clientes vão junto. Nada fica registrado."),
    ("💤", "Cliente some e ninguém percebe", "Sem alerta de inatividade, a recompra simplesmente não acontece."),
    ("📦", "Estoque e vendas desconectados", "Grade por cor/tamanho no caderno; encalhe e ruptura passam batido."),
    ("📣", "Disparo genérico no WhatsApp", "Mesma mensagem para todo mundo — baixa conversão e descadastro."),
]
x = 0.9; y = 2.05; w = 5.75; h = 1.95
for i, (ic, t, d) in enumerate(dores):
    cx = x + (i % 2) * (w + 0.35)
    cy = y + (i // 2) * (h + 0.3)
    rect(s, cx, cy, w, h)
    tf = box(s, cx + 0.3, cy + 0.25, w - 0.6, h - 0.4)
    text(tf, f"{ic}  {t}", size=19, bold=True, color=INK, font=SERIF, first=True, space_after=8)
    text(tf, d, size=15, color=INK_SOFT, line=1.1)

# ───────────────────────── 3. Visão geral ─────────────────────────
s = slide()
kicker(s, "O que é")
title(s, "Um sistema, toda a operação de moda")
mods = [
    ("📊", "Dashboard & metas"), ("🛒", "Vendas + PDV"), ("🏷️", "Estoque em grade"),
    ("👗", "Carteira inteligente"), ("🤖", "IA Vendedora"), ("📱", "WhatsApp por vendedora"),
    ("🎯", "Disparos personalizados"), ("🔁", "Recuperação de inativos"), ("🏆", "Ranking & comissão"),
    ("🛍️", "Catálogo da vendedora"), ("🧥", "Provador virtual IA"), ("📦", "Atacado & portal cliente"),
]
cols = 3; cw = 3.78; ch = 1.16; gx = 0.2; gy = 0.22; ox = 0.9; oy = 2.0
for i, (ic, name) in enumerate(mods):
    cx = ox + (i % cols) * (cw + gx)
    cy = oy + (i // cols) * (ch + gy)
    rect(s, cx, cy, cw, ch)
    tf = box(s, cx + 0.25, cy + 0.0, cw - 0.4, ch)
    tf.word_wrap = True
    p = text(tf, f"{ic}  {name}", size=16, bold=True, color=INK, first=True)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
text(box(s, 0.9, 6.85, 11.5, 0.5),
     [("★ Diferencial: ", {"color": ACCENT, "bold": True}),
      ("Radar de Oportunidades — cruza estoque parado × perfil do cliente e dispara em 1 clique.", {"color": INK_SOFT})],
     size=14, first=True)

# ───────────────────────── 4. Arquitetura ─────────────────────────
s = slide()
kicker(s, "Como é construído")
title(s, "SaaS multi-tenant, isolado por hierarquia")
# hierarquia
nodes = ["Rede / Marca", "Loja", "Equipe", "Vendedora", "Cliente"]
nx = 0.9; ny = 2.2; nw = 2.15; nh = 0.95
for i, n in enumerate(nodes):
    cx = nx + i * (nw + 0.18)
    fill = ACCENT if i == 0 else WHITE
    col = WHITE if i == 0 else INK
    rect(s, cx, ny, nw, nh, fill=fill)
    tf = box(s, cx, ny, nw, nh); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(tf, n, size=15, bold=True, color=col, align=PP_ALIGN.CENTER, first=True)
    if i < len(nodes) - 1:
        text(box(s, cx + nw - 0.04, ny + 0.18, 0.3, 0.6), "›", size=26, color=ACCENT, bold=True, first=True)
# pilares
pil = [
    ("Stack", "Node.js · Fastify · Prisma · PostgreSQL · React + Vite"),
    ("Segurança", "JWT + RBAC; isolamento de carteira aplicado no backend, não só na UI"),
    ("Integrações", "WhatsApp (Evolution API) · IA (Claude) · Instagram/Meta Ads"),
    ("Planos SaaS", "Start · Pro · Elite (multi-lojas) com limites aplicados no servidor"),
]
py = 3.7
for i, (t, d) in enumerate(pil):
    cx = 0.9 + (i % 2) * 5.95
    cy = py + (i // 2) * 1.45
    rect(s, cx, cy, 5.7, 1.25)
    tf = box(s, cx + 0.3, cy + 0.2, 5.1, 0.95)
    text(tf, t, size=16, bold=True, color=ACCENT, first=True, space_after=4)
    text(tf, d, size=13.5, color=INK_SOFT, line=1.08)

# ───────────────────────── 5. Status / roadmap ─────────────────────────
s = slide()
kicker(s, "Onde estamos")
title(s, "7 fases entregues e rodando")
done = [
    ("Fase 1 — Fundação", "Auth multi-tenant, RBAC, CRUD com grade (cor × tamanho)."),
    ("Fase 2 — Vendas + Estoque", "PDV com baixa atômica + dashboards hierárquicos."),
    ("Fase 3 — CRM", "Carteira por vendedora + segmentação automática."),
    ("Fase 4 — WhatsApp", "Disparos por segmento (IA) + réguas de inativos + webhook."),
    ("Fase 5 — IA / Radar ★", "Estoque inteligente + Radar de Oportunidades (1 clique)."),
    ("Fase 6 — Comissão + Gamificação", "Comissão encadeada, ranking, mural."),
    ("Fase 7 — Expansão", "Provador Virtual IA + Atacado."),
]
y = 1.78
for t, d in done:
    rect(s, 0.9, y, 11.5, 0.62, fill=WHITE, line=CARD_BD)
    rect(s, 0.9, y, 0.13, 0.62, fill=OK, line=None, round=False)
    tf = box(s, 1.25, y + 0.07, 11, 0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(tf, [("✓  ", {"color": OK, "bold": True, "size": 14}), (t, {"bold": True, "size": 14.5, "font": SERIF}),
              ("   —   " + d, {"size": 12.5, "color": INK_SOFT})], first=True)
    y += 0.74
text(box(s, 0.9, y + 0.04, 11.5, 0.5),
     [("+ extras: ", {"bold": True, "color": ACCENT}),
      ("forma de recebimento · produto/moda · gestor de estoque+entradas · transferências (radar de extravio) · gestão de gestores de estoque.", {"color": INK_SOFT})],
     size=12.5, first=True)

# ───────────────────────── Divisória DEMO ─────────────────────────
s = slide(INK)
rect(s, 0, 3.05, SW.inches, 1.4, fill=ACCENT, line=None, round=False)
text(box(s, 0, 3.2, 13.333, 1.1), "DEMONSTRAÇÃO AO VIVO", size=40, color=WHITE,
     font=SERIF, bold=True, align=PP_ALIGN.CENTER, first=True)
text(box(s, 0, 5.0, 13.333, 0.6), "5 cenas · do gerente à vendedora à rede",
     size=18, color=RGBColor(0xD8,0xCF,0xC8), align=PP_ALIGN.CENTER, first=True)

# ───────────────────────── Acessos de teste ─────────────────────────
s = slide()
kicker(s, "Logins para a demonstração")
title(s, "Acessos de teste")
acessos = [
    ("Super Admin · SaaS", "admin@modacrm.com.br", "admin123", False),
    ("Gestor · rede Luna Brand (ELITE)", "gestor@lunabrand.com.br", "gestor123", True),
    ("Gestor de Estoque · rede", "estoquista@lunabrand.com.br", "estoque123", False),
    ("Gerente · Loja Demo Moda", "maria@lojademo.com.br", "demo123", True),
    ("Gerente · Loja Shopping", "paula@lojashopping.com.br", "demo123", False),
    ("Vendedora · Demo Moda", "camila@lojademo.com.br", "demo123", True),
    ("Vendedora · Demo Moda", "ana@lojademo.com.br", "demo123", False),
    ("Vendedora · Loja Shopping", "julia@lojashopping.com.br", "demo123", False),
]
cx, cw = 0.9, 11.5
# cabeçalho
ty = 2.0
hdr = rect(s, cx, ty, cw, 0.5, fill=INK, line=None)
for label, lx, lw in [("Papel", 0.35, 4.6), ("E-mail", 5.1, 4.4), ("Senha", 9.6, 2.6)]:
    tfh = box(s, cx + lx, ty + 0.04, lw, 0.42); tfh.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(tfh, label, size=13, bold=True, color=WHITE, first=True)
ty += 0.5
for papel, email, senha, dest in acessos:
    rh = 0.6
    rect(s, cx, ty, cw, rh, fill=(ACCENT_SFT if dest else WHITE), line=CARD_BD)
    if dest:
        rect(s, cx, ty, 0.1, rh, fill=ACCENT, line=None, round=False)
    star = "★ " if dest else ""
    t1 = box(s, cx + 0.35, ty, 4.6, rh); t1.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(t1, [(star, {"color": ACCENT, "bold": True}), (papel, {"bold": dest})], size=13.5, color=INK, first=True)
    t2 = box(s, cx + 5.1, ty, 4.4, rh); t2.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(t2, email, size=13.5, color=INK, first=True)
    t3 = box(s, cx + 9.6, ty, 2.6, rh); t3.vertical_anchor = MSO_ANCHOR.MIDDLE
    text(t3, senha, size=13.5, color=ACCENT, bold=True, font="Consolas", first=True)
    ty += rh
text(box(s, 0.9, ty + 0.18, 11.5, 0.5),
     [("★ usados na demo  ·  ", {"color": ACCENT, "bold": True}),
      ("abra http://localhost:5184  ·  sequência: Maria → Camila → Gestor", {"color": INK_SOFT})],
     size=13.5, first=True)


def demo_slide(num, total, titulo, persona, login, passos, destaque):
    s = slide()
    kicker(s, f"Demo {num} de {total}")
    title(s, titulo, size=30)
    # passos
    tf = box(s, 0.9, 2.15, 7.2, 4.6)
    first = True
    for p in passos:
        text(tf, [("●  ", {"color": ACCENT, "bold": True}), (p, {})],
             size=17, color=INK, line=1.12, space_after=12, first=first)
        first = False
    # painel lateral persona + destaque
    rect(s, 8.5, 2.15, 3.9, 1.5, fill=ACCENT_SFT, line=None)
    tf = box(s, 8.75, 2.32, 3.45, 1.2)
    text(tf, "PERSONA", size=12, color=ACCENT, bold=True, first=True, space_after=3)
    text(tf, persona, size=16, bold=True, color=INK, first=False, space_after=2)
    text(tf, login, size=12.5, color=INK_SOFT)
    rect(s, 8.5, 3.85, 3.9, 2.9, fill=INK, line=None)
    tf = box(s, 8.75, 4.05, 3.45, 2.55)
    text(tf, "💡 PONTO-CHAVE", size=12, color=RGBColor(0xE8,0xA8,0x7C), bold=True, first=True, space_after=6)
    text(tf, destaque, size=15.5, color=WHITE, line=1.18)
    return s

# ── Demo 1 ──
demo_slide(1, 5, "Dashboard do gerente", "Maria · Gerente", "maria@lojademo.com.br",
    ["Faturamento do dia e do mês, ticket médio",
     "Ranking de vendedoras com % da meta",
     "Produtos mais vendidos e melhores clientes",
     "Estoque crítico destacado automaticamente"],
    "Em uma tela, a gerente vê a saúde da loja e quem está perto (ou longe) da meta.")

# ── Demo 2 ──
demo_slide(2, 5, "Vender e baixar estoque", "Maria · Gerente", "maria@lojademo.com.br",
    ["Nova venda: cliente, vendedora e itens",
     "Variação por cor/tamanho com estoque à vista",
     "Preço varejo ou atacado, desconto e total ao vivo",
     "Ao salvar, o estoque baixa por SKU na hora"],
    "Venda e estoque deixam de ser dois mundos: cada venda movimenta o SKU e alimenta o dashboard.")

# ── Demo 3 (o wow) ──
demo_slide(3, 5, "Segmentação automática (CRM)", "Maria · Gerente", "maria@lojademo.com.br",
    ["Carteira mostra a distribuição por segmento",
     "Clique em “Recalcular segmentação”",
     "Clientes se reclassificam pelo histórico real",
     "Abra uma ficha: histórico + alerta de recuperação"],
    "O sistema lê as compras e decide sozinho quem é VIP, Frequente, Novo, Atacado ou Inativo — e aponta quem recuperar.")

# ── Demo 4 ──
demo_slide(4, 5, "Isolamento da carteira", "Camila · Vendedora", "camila@lojademo.com.br",
    ["Login como vendedora",
     "Vê apenas os próprios clientes",
     "Dashboard mostra só o desempenho dela",
     "Não acessa vendas/carteira das colegas"],
    "Privacidade da carteira garantida no backend — não é só esconder na tela. Confiança da equipe.")

# ── Demo 5 ──
demo_slide(5, 5, "Visão da rede (multi-lojas)", "Gustavo · Gestor", "gestor@lunabrand.com.br",
    ["Login como gestor da marca (plano Elite)",
     "Consolidado de faturamento da rede",
     "Comparativo loja a loja",
     "Drill-down para o dashboard de cada loja"],
    "Do balcão à rede: a marca enxerga todas as lojas em um só painel, com isolamento por loja preservado.")

# ───────────────────────── Diferencial: Radar ─────────────────────────
s = slide(INK)
kicker(s, "O diferencial", x=0.9, y=0.7)
text(box(s, 0.9, 0.95, 11.5, 1.0), "★ Radar de Oportunidades", size=34, color=WHITE,
     font=SERIF, bold=True, first=True)
bar(s, y=1.75, x=0.9)
steps = [
    ("1", "Estoque parado", "A IA identifica peças encalhadas e em risco de ruptura."),
    ("2", "Perfil do cliente", "Cruza com quem compra aquele estilo, cor, tamanho ou marca."),
    ("3", "Disparo em 1 clique", "Gera a campanha de WhatsApp personalizada e envia pela vendedora dona."),
]
x = 0.9; y = 2.3; w = 3.78
for i, (n, t, d) in enumerate(steps):
    cx = x + i * (w + 0.35)
    rect(s, cx, y, w, 2.7, fill=RGBColor(0x39,0x2F,0x2B), line=None)
    tf = box(s, cx + 0.35, y + 0.3, w - 0.7, 2.2)
    text(tf, n, size=30, color=ACCENT, bold=True, font=SERIF, first=True, space_after=6)
    text(tf, t, size=19, color=WHITE, bold=True, space_after=8)
    text(tf, d, size=14.5, color=RGBColor(0xD8,0xCF,0xC8), line=1.18)
text(box(s, 0.9, 5.5, 11.5, 0.8),
     "Transforma estoque parado em receita — a oportunidade certa, para o cliente certo, pela vendedora certa.",
     size=16, color=RGBColor(0xE8,0xA8,0x7C), first=True)

# ───────────────────────── Modelo de receita ─────────────────────────
s = slide()
kicker(s, "Modelo de negócio")
title(s, "SaaS recorrente, por porte da operação")
planos = [
    ("Start", "R$ 97", "/mês", ["Até 2 vendedoras", "1 loja", "CRM + Vendas + Estoque"], False),
    ("Pro", "R$ 297", "/mês", ["Até 10 vendedoras", "1 loja", "WhatsApp + disparos + IA básica"], True),
    ("Elite", "R$ 697", "/mês", ["Vendedoras ilimitadas", "Multi-lojas + rede", "IA avançada + Radar + Meta Ads"], False),
]
x = 0.9; y = 2.1; w = 3.78; h = 4.2
for i, (nome, preco, per, feats, dest) in enumerate(planos):
    cx = x + i * (w + 0.35)
    rect(s, cx, y, w, h, fill=(ACCENT if dest else WHITE), line=(None if dest else CARD_BD))
    col = WHITE if dest else INK
    sub = RGBColor(0xF5,0xE3,0xDA) if dest else INK_SOFT
    tf = box(s, cx + 0.4, y + 0.35, w - 0.8, h - 0.7)
    if dest:
        text(tf, "MAIS POPULAR", size=11, color=WHITE, bold=True, first=True, space_after=4)
        text(tf, nome, size=26, color=col, bold=True, font=SERIF, space_after=2)
    else:
        text(tf, nome, size=26, color=col, bold=True, font=SERIF, first=True, space_after=2)
    text(tf, [(preco, {"size": 30, "bold": True, "color": col}), (per, {"size": 15, "color": sub})], space_after=12)
    for f in feats:
        text(tf, [("✓  ", {"color": (WHITE if dest else OK), "bold": True}), (f, {"color": col})], size=14.5, line=1.1, space_after=8)

# ───────────────────────── Encerramento ─────────────────────────
s = slide(INK)
rect(s, 0, 0, 0.35, SH.inches, fill=ACCENT, line=None, round=False)
tf = box(s, 1.1, 2.5, 11, 2.6)
text(tf, "Vamos vestir sua operação de inteligência.", size=36, color=WHITE,
     font=SERIF, bold=True, first=True, space_after=14, line=1.05)
text(tf, "ModaCRM AI — vendas, relacionamento, estoque e IA em um só lugar.",
     size=19, color=RGBColor(0xD8,0xCF,0xC8), space_after=20)
text(tf, [("Próximo passo: ", {"color": RGBColor(0xE8,0xA8,0x7C), "bold": True}),
          ("Fase 8 — Billing SaaS (planos Start/Pro/Elite com cobrança recorrente).", {"color": RGBColor(0xA8,0x9D,0x96)})],
     size=16)

# ── Rodapé: marca + numeração (pula a capa) ──
total_slides = len(SLIDES)
for idx, (sl, dark) in enumerate(SLIDES, start=1):
    if idx == 1:
        continue
    col = RGBColor(0x9A, 0x8F, 0x88) if dark else INK_SOFT
    tl = box(sl, 0.9, 7.06, 4.0, 0.3)
    text(tl, "ModaCRM AI", size=10, color=col, first=True)
    tr = box(sl, 11.4, 7.06, 1.5, 0.3)
    text(tr, f"{idx} / {total_slides}", size=10, color=col, align=PP_ALIGN.RIGHT, first=True)

import os
out = os.path.join(os.path.dirname(__file__), "ModaCRM-AI.pptx")
prs.save(out)
print("OK:", out, "| slides:", len(prs.slides._sldIdLst))
